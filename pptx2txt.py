import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.exc import PackageNotFoundError
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape

version = "2025-10-28"


def log(message: str) -> None:
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    print(f"{timestamp} {message}", file=sys.stderr)


def extract_lines_from_shape(shape: BaseShape) -> list[str]:
    shape_lines = []
    if isinstance(shape, Shape):
        for paragraph in shape.text_frame.paragraphs:
            stripped = paragraph.text.strip()
            if stripped:
                shape_lines.append(stripped)
    elif isinstance(shape, GraphicFrame) and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                stripped = cell.text.strip()
                if stripped:
                    shape_lines.append(stripped)
    elif isinstance(shape, GroupShape):
        for child_shape in shape.shapes:
            shape_lines += extract_lines_from_shape(child_shape)
    return shape_lines


log(f"pptx2txt - version {version} by Shinichi Akiyama")

if len(sys.argv) < 2:
    log("ERROR: No input files specified")
    log("Usage: pptx2txt <file1.pptx> [file2.pptx] ...")
    sys.exit(1)

for pptx_file in sys.argv[1:]:
    pptx_path = Path(pptx_file)
    if not pptx_path.exists():
        log(f"ERROR: File not found: {pptx_file}")
        sys.exit(1)

    try:
        presentation = Presentation(pptx_file)
        log(f"{pptx_file} was opened.")
    except PackageNotFoundError:
        log(f"ERROR: Invalid PPTX file: {pptx_file}")
        sys.exit(1)
    except Exception as e:
        log(f"ERROR: Failed to open {pptx_file}: {e}")
        sys.exit(1)

    lines = []
    for slide_index, slide in enumerate(presentation.slides):
        lines.append(f"--- Slide {slide_index + 1} ---")
        for shape in slide.shapes:
            lines += extract_lines_from_shape(shape)

    text = "\n".join(lines) + "\n"
    text_path = pptx_path.with_suffix(".txt")
    try:
        text_path.write_text(text, encoding="utf-8")
        log(f"{text_path} was saved.")
    except OSError as e:
        log(f"ERROR: Cannot write to {text_path}: {e}")
        sys.exit(1)

log("All files were processed.")
